import json
from pathlib import Path
import argparse

from PIL import Image
import cv2, numpy as np

# PPTX
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR

# DOCX (optional; PPTX is the reliable editor)
from docx import Document
from docx.shared import Inches as DocxInches, Emu
from docx.oxml import OxmlElement
from docx.oxml.ns  import qn, nsmap

# OCR
import pytesseract

# Ensure Wordprocessing Shape namespace for DOCX textboxes
nsmap.setdefault('wps', 'http://schemas.microsoft.com/office/word/2010/wordprocessingShape')

EMU_PER_INCH = 914400
def px_to_inches(px, dpi=96.0): return px / dpi
def px_to_emu(px, dpi=96.0):    return int(px_to_inches(px, dpi) * EMU_PER_INCH)

# --------- background erase (inpaint) ----------
def inpaint_background(img_path, boxes, expand_px=2, radius=3):
    img = cv2.imread(img_path, cv2.IMREAD_COLOR)
    if img is None: raise RuntimeError(f"Could not read {img_path}")
    h, w = img.shape[:2]
    mask = np.zeros((h, w), dtype=np.uint8)
    for b in boxes:
        x0 = max(0, b["left"] - expand_px)
        y0 = max(0, b["top"]  - expand_px)
        x1 = min(w, b["left"] + b["width"]  + expand_px)
        y1 = min(h, b["top"]  + b["height"] + expand_px)
        cv2.rectangle(mask, (x0,y0), (x1,y1), 255, thickness=-1)
    mask = cv2.dilate(mask, np.ones((3,3), np.uint8), iterations=1)
    cleaned = cv2.inpaint(img, mask, radius, cv2.INPAINT_TELEA)
    out_path = str(Path(img_path).with_suffix("")) + "_clean.png"
    cv2.imwrite(out_path, cleaned)
    return out_path

# --------- OCR ----------
def ocr_crop(pil_img, lang="eng", psm=6):
    cfg = f"--psm {psm}"
    txt = pytesseract.image_to_string(pil_img, lang=lang, config=cfg)
    return (txt or "").strip()

# --------- PPTX ----------
def add_pptx_bg(slide, img_path, w_in, h_in):
    slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(w_in), height=Inches(h_in))

def add_pptx_text(
    slide, boxes, dpi, font_name=None, font_size_pt=None, rtl=False,
    ppt_autofit=True, ppt_margin_pt=1.5, debug_outline=False
):
    """
    Places text strictly inside each user-drawn box:
      - zero (or tiny) margins
      - word-wrap on
      - shrink-to-fit (if ppt_autofit=True)
      - tight paragraph spacing
    """
    margin_in = max(0.0, float(ppt_margin_pt)) / 72.0
    for b in boxes:
        li = px_to_inches(b['left'],  dpi)
        ti = px_to_inches(b['top'],   dpi)
        wi = px_to_inches(b['width'], dpi)
        hi = px_to_inches(b['height'],dpi)

        if debug_outline:
            rect = slide.shapes.add_shape(1, Inches(li), Inches(ti), Inches(wi), Inches(hi))
            rect.line.width = Pt(0.5)

        tb = slide.shapes.add_textbox(Inches(li), Inches(ti), Inches(wi), Inches(hi))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        # margins
        tf.margin_left   = Inches(margin_in)
        tf.margin_right  = Inches(margin_in)
        tf.margin_top    = Inches(margin_in)
        tf.margin_bottom = Inches(margin_in)

        if ppt_autofit:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # PPT will shrink to fit

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        if rtl: p.rtl = True
        p.space_before = Pt(0)
        p.space_after  = Pt(0)
        p.line_spacing = 1.05

        run = p.add_run()
        run.text = b.get("text", "[edit]")
        f = run.font
        if font_name: f.name = font_name
        if font_size_pt:
            f.size = Pt(font_size_pt)   # you can force a size
        else:
            # start large; autofit will shrink it down to perfect fit
            f.size = Pt(80) if ppt_autofit else Pt(12)

# --------- DOCX (optional; PPTX recommended) ----------
def add_docx_bg_header(doc, img_path):
    section = doc.sections[-1]
    p = section.header.paragraphs[0] if section.header.paragraphs else section.header.add_paragraph()
    p.add_run().add_picture(img_path, width=section.page_width, height=section.page_height)

def add_docx_textbox(paragraph, left_emu, top_emu, width_emu, height_emu,
                     text, font_name=None, font_size_pt=None, rtl=False):
    drawing = OxmlElement('w:drawing')
    anchor = OxmlElement('wp:anchor')
    for k in ('distT','distB','distL','distR'): anchor.set(k,'0')
    anchor.set('simplePos','0'); anchor.set('relativeHeight','251658240')
    anchor.set('behindDoc','0'); anchor.set('locked','0'); anchor.set('layoutInCell','1'); anchor.set('allowOverlap','1')
    sp = OxmlElement('wp:simplePos'); sp.set('x','0'); sp.set('y','0'); anchor.append(sp)
    ph = OxmlElement('wp:positionH'); ph.set('relativeFrom','page')
    offH = OxmlElement('wp:posOffset'); offH.text = str(left_emu); ph.append(offH); anchor.append(ph)
    pv = OxmlElement('wp:positionV'); pv.set('relativeFrom','page')
    offV = OxmlElement('wp:posOffset'); offV.text = str(top_emu);  pv.append(offV); anchor.append(pv)
    ext = OxmlElement('wp:extent'); ext.set('cx', str(width_emu)); ext.set('cy', str(height_emu)); anchor.append(ext)
    ee = OxmlElement('wp:effectExtent'); [ee.set(k,'0') for k in ('l','t','r','b')]
    anchor.append(ee)
    anchor.append(OxmlElement('wp:wrapNone'))

    graphic = OxmlElement('a:graphic')
    gd = OxmlElement('a:graphicData'); gd.set('uri','http://schemas.microsoft.com/office/word/2010/wordprocessingShape')
    wsp = OxmlElement('wps:wsp')

    spPr = OxmlElement('a:spPr')
    xfrm = OxmlElement('a:xfrm'); ext2 = OxmlElement('a:ext'); ext2.set('cx', str(width_emu)); ext2.set('cy', str(height_emu))
    xfrm.append(ext2); spPr.append(xfrm)
    geom = OxmlElement('a:prstGeom'); geom.set('prst','rect'); geom.append(OxmlElement('a:avLst')); spPr.append(geom)
    spPr.append(OxmlElement('a:noFill'))
    ln = OxmlElement('a:ln'); ln.append(OxmlElement('a:noFill')); spPr.append(ln)
    wsp.append(spPr)

    txbx = OxmlElement('wps:txbx'); txbxContent = OxmlElement('w:txbxContent')
    wp = OxmlElement('w:p')
    if rtl:
        pPr = OxmlElement('w:pPr'); bidi = OxmlElement('w:bidi'); bidi.set(qn('w:val'),'on'); pPr.append(bidi); wp.append(pPr)
    wr = OxmlElement('w:r'); rPr = OxmlElement('w:rPr')
    if font_size_pt:
        sz = OxmlElement('w:sz'); sz.set(qn('w:val'), str(int(font_size_pt*2))); rPr.append(sz)
    if font_name:
        rFonts = OxmlElement('w:rFonts')
        for k in ('w:ascii','w:hAnsi','w:eastAsia','w:cs'): rFonts.set(qn(k), font_name)
        rPr.append(rFonts)
    wr.append(rPr)
    wt = OxmlElement('w:t'); wt.text = text; wr.append(wt); wp.append(wr)
    txbxContent.append(wp); txbx.append(txbxContent)

    bodyPr = OxmlElement('wps:bodyPr')
    for k in ('lIns','tIns','rIns','bIns'): bodyPr.set(k,'0')
    bodyPr.set('wrap','square')
    wsp.append(txbx); wsp.append(bodyPr)

    gd.append(wsp); graphic.append(gd); anchor.append(graphic)
    drawing.append(anchor)
    paragraph.add_run()._r.append(drawing)

# --------- builder ----------
def build_from_json(json_path, out_pptx=None, out_docx=None, dpi=300.0,
                    font_name=None, font_size_pt=None, rtl=False,
                    erase=False, expand_px=2, radius=3,
                    ocr_prefill=False, lang="eng", tesseract_path=None, psm=6,
                    ppt_autofit=True, ppt_margin_pt=1.5, debug_outline=False):
    if tesseract_path:
        pytesseract.pytesseract.tesseract_cmd = tesseract_path

    meta = json.loads(Path(json_path).read_text(encoding="utf-8"))
    images = meta["images"]

    prs = Presentation() if out_pptx else None
    doc = Document()     if out_docx else None

    for i, imeta in enumerate(images):
        img_path = imeta["path"]
        im = Image.open(img_path).convert("RGB")
        w_px, h_px = im.size

        # OCR prefill
        boxes = []
        for b in imeta["boxes"]:
            bb = dict(b)
            if ocr_prefill:
                crop = im.crop((b["left"], b["top"], b["left"]+b["width"], b["top"]+b["height"]))
                txt = ocr_crop(crop, lang=lang, psm=psm)
                bb["text"] = txt if txt else "[edit]"
            else:
                bb["text"] = "[edit]"
            boxes.append(bb)

        # erase underlying text if requested
        bg_path = inpaint_background(img_path, boxes, expand_px=expand_px, radius=radius) if erase else img_path

        # sizes
        w_in = px_to_inches(w_px, dpi); h_in = px_to_inches(h_px, dpi)

        if prs:
            prs.slide_width = Inches(w_in); prs.slide_height = Inches(h_in)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_pptx_bg(slide, bg_path, w_in, h_in)
            add_pptx_text(
                slide, boxes, dpi, font_name, font_size_pt, rtl=rtl,
                ppt_autofit=ppt_autofit, ppt_margin_pt=ppt_margin_pt, debug_outline=debug_outline
            )

        if doc:
            if i > 0: doc.add_page_break()
            section = doc.sections[-1]
            section.page_width  = Emu(w_in * EMU_PER_INCH)
            section.page_height = Emu(h_in * EMU_PER_INCH)
            m = DocxInches(0.25)
            section.left_margin = m; section.right_margin = m; section.top_margin = m; section.bottom_margin = m
            add_docx_bg_header(doc, bg_path)
            for b in boxes:
                p = doc.add_paragraph()
                add_docx_textbox(
                    p,
                    left_emu=px_to_emu(b["left"], dpi),
                    top_emu=px_to_emu(b["top"], dpi),
                    width_emu=px_to_emu(b["width"], dpi),
                    height_emu=px_to_emu(b["height"], dpi),
                    text=b["text"],
                    font_name=font_name, font_size_pt=font_size_pt,
                    rtl=rtl
                )

    if prs: prs.save(out_pptx); print("Saved", out_pptx)
    if doc: doc.save(out_docx); print("Saved", out_docx)

if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--json", required=True)
    ap.add_argument("--out_pptx", default=None)
    ap.add_argument("--out_docx", default=None)
    ap.add_argument("--dpi", type=float, default=300.0)
    ap.add_argument("--font", default=None)
    ap.add_argument("--size", type=float, default=None, help="If omitted, shrink-to-fit handles sizing when --ppt_autofit")
    ap.add_argument("--rtl", action="store_true")
    ap.add_argument("--erase", action="store_true")
    ap.add_argument("--expand_px", type=int, default=2)
    ap.add_argument("--radius", type=int, default=3)
    ap.add_argument("--ocr_prefill", action="store_true")
    ap.add_argument("--lang", default="eng")
    ap.add_argument("--tesseract_path", default=None)
    ap.add_argument("--psm", type=int, default=6)
    ap.add_argument("--ppt_autofit", action="store_true", help="Shrink text to fit each textbox")
    ap.add_argument("--ppt_margin_pt", type=float, default=1.5, help="Textbox margins in points")
    ap.add_argument("--debug_outline", action="store_true")
    args = ap.parse_args()

    build_from_json(
        args.json,
        out_pptx=args.out_pptx,
        out_docx=args.out_docx,
        dpi=args.dpi,
        font_name=args.font,
        font_size_pt=args.size,
        rtl=args.rtl,
        erase=args.erase,
        expand_px=args.expand_px,
        radius=args.radius,
        ocr_prefill=args.ocr_prefill,
        lang=args.lang,
        tesseract_path=args.tesseract_path,
        psm=args.psm,
        ppt_autofit=args.ppt_autofit,
        ppt_margin_pt=args.ppt_margin_pt,
        debug_outline=args.debug_outline,
    )
